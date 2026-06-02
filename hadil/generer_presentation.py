# -*- coding: utf-8 -*-
"""
Dind'Or PFE - Generateur de la presentation PowerPoint
=======================================================

Utilisation :
    1. (une seule fois)  pip install python-pptx
    2.                   python generer_presentation.py
    3. Le fichier "Dindor_PFE_Presentation.pptx" sera cree dans ce dossier.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

NAVY=RGBColor(0x0F,0x1B,0x2D); NAVY_D=RGBColor(0x08,0x12,0x20)
GOLD=RGBColor(0xD4,0xA4,0x37); GOLD_S=RGBColor(0xE9,0xC4,0x6A)
CREAM=RGBColor(0xFA,0xF6,0xEE); WHITE=RGBColor(0xFF,0xFF,0xFF)
SLATE=RGBColor(0x4A,0x55,0x68); DARK=RGBColor(0x1A,0x23,0x32)
RED=RGBColor(0xC0,0x4A,0x3D)

W,H=13.333,7.5
prs=Presentation(); prs.slide_width=Inches(W); prs.slide_height=Inches(H)
B=prs.slide_layouts[6]

def r(s,x,y,w,h,f):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=f; sh.line.fill.background(); sh.shadow.inherit=False
    return sh
def rr(s,x,y,w,h,f,rd=0.05):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.adjustments[0]=rd; sh.fill.solid(); sh.fill.fore_color.rgb=f
    sh.line.fill.background(); sh.shadow.inherit=False
    return sh
def o(s,x,y,d,f):
    sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(d),Inches(d))
    sh.fill.solid(); sh.fill.fore_color.rgb=f; sh.line.fill.background(); sh.shadow.inherit=False
    return sh
def t(s,x,y,w,h,txt,sz=14,c=DARK,bd=False,a=PP_ALIGN.LEFT,it=False):
    b=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=b.text_frame; tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=0
    tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.TOP
    p=tf.paragraphs[0]; p.alignment=a; rn=p.add_run(); rn.text=txt
    rn.font.name="Calibri"; rn.font.size=Pt(sz); rn.font.bold=bd
    rn.font.italic=it; rn.font.color.rgb=c
    return b
def pg(s,n):
    t(s,W-1.0,H-0.45,0.8,0.3,f"{n} / 10",sz=10,c=SLATE,a=PP_ALIGN.RIGHT,it=True)
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
    sh.fill.solid(); sh.fill.fore_color.rgb=c
    sh.line.fill.background(); sh.shadow.inherit=False
    sp=sh._element.getparent(); sp.remove(sh._element); sp.insert(2,sh._element)
def hdr(s,lbl,ttl,dk=False):
    t(s,0.6,0.55,8,0.32,lbl.upper(),sz=11,c=GOLD,bd=True)
    r(s,0.6,0.93,0.4,0.06,GOLD)
    t(s,0.6,1.05,12.5,0.85,ttl,sz=32,bd=True,c=(WHITE if dk else DARK))

# ===== SLIDE 1 - COUVERTURE =====
s=prs.slides.add_slide(B); bg(s,NAVY)
big=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(9.5),Inches(-2),Inches(7),Inches(7))
big.fill.solid(); big.fill.fore_color.rgb=GOLD; big.line.fill.background(); big.shadow.inherit=False
sp=big.fill.fore_color._xFill
al=etree.SubElement(sp.find(qn('a:srgbClr')),qn('a:alpha')); al.set('val','15000')
r(s,0.7,0.7,0.35,0.35,GOLD)
t(s,1.2,0.65,8,0.4,"PROJET DE FIN D'ETUDES  |  2025-2026",sz=12,c=GOLD_S,bd=True)
t(s,0.7,2.0,11,1.0,"Dind'Or ERP",sz=64,bd=True,c=WHITE)
t(s,0.7,3.0,11,0.7,"Augmente par l'Intelligence Artificielle",sz=26,c=GOLD,it=True)
r(s,0.7,4.0,0.7,0.05,GOLD)
t(s,0.7,4.2,11,0.5,"Systeme de gestion integre + Prediction des prix d'achat",sz=18,c=CREAM)
cy=5.7
r(s,0.7,cy,5.5,0.05,GOLD)
t(s,0.7,cy+0.15,6,0.35,"ENTREPRISE D'ACCUEIL",sz=10,c=GOLD,bd=True)
t(s,0.7,cy+0.5,6,0.4,"Dind'Or - Monastir, Tunisie",sz=16,c=WHITE,bd=True)
t(s,0.7,cy+0.95,6.5,0.4,"Elevage et vente de volaille (depuis 2014)",sz=12,c=CREAM,it=True)
t(s,7.5,cy+0.15,5,0.35,"SOUTENANCE TECHNIQUE",sz=10,c=GOLD,bd=True,a=PP_ALIGN.RIGHT)
t(s,7.5,cy+0.5,5,0.4,"Presentation du projet",sz=16,c=WHITE,bd=True,a=PP_ALIGN.RIGHT)
t(s,7.5,cy+0.95,5,0.4,"Spring Boot / Angular / FastAPI",sz=12,c=CREAM,it=True,a=PP_ALIGN.RIGHT)

# ===== SLIDE 2 - PLAN =====
s=prs.slides.add_slide(B); bg(s,CREAM); hdr(s,"Sommaire","Plan de la presentation")
secs=[("01","Presentation du projet","Entreprise et contexte"),
      ("02","Objectifs","Generaux et specifiques"),
      ("03","Problematique","Le defi metier"),
      ("04","Solution proposee","ERP + IA"),
      ("05","Architecture","3 couches + microservice"),
      ("06","Environnement de travail","Stack technologique"),
      ("07","Methodologie SCRUM","Sprints et organisation"),
      ("08","Conclusion","Bilan et perspectives")]
cw,rh=5.8,0.85
for i,(n,tt,sb) in enumerate(secs):
    cx=0.6+(i%2)*(cw+0.3); cy=2.2+(i//2)*(rh+0.3)
    rr(s,cx,cy,cw,rh,WHITE,0.15); r(s,cx,cy,0.08,rh,GOLD)
    t(s,cx+0.3,cy+0.12,0.8,0.6,n,sz=22,bd=True,c=GOLD)
    t(s,cx+1.2,cy+0.13,cw-1.3,0.4,tt,sz=15,bd=True,c=DARK)
    t(s,cx+1.2,cy+0.48,cw-1.3,0.35,sb,sz=11,c=SLATE,it=True)
pg(s,2)

# ===== SLIDE 3 - PRESENTATION =====
s=prs.slides.add_slide(B); bg(s,CREAM)
hdr(s,"01 | Presentation","Une entreprise tunisienne, un projet ambitieux")
lx,ly=0.6,2.2
rr(s,lx,ly,5.8,4.5,WHITE); r(s,lx,ly,5.8,0.6,NAVY)
t(s,lx+0.3,ly+0.12,5.2,0.4,"L'ENTREPRISE - DIND'OR",sz=14,bd=True,c=GOLD)
fc=[("Secteur","Elevage et vente de volaille"),("Fondation","2014"),
    ("Siege","Monastir, Tunisie"),("Activites","Achat en gros, B2B/B2C"),
    ("Defi metier","Volatilite des prix d'achat")]
fy=ly+0.85
for k,v in fc:
    o(s,lx+0.3,fy+0.1,0.12,GOLD)
    t(s,lx+0.55,fy,1.7,0.4,k,sz=11,bd=True,c=SLATE)
    t(s,lx+2.3,fy,3.4,0.4,v,sz=12,c=DARK); fy+=0.65
rx=lx+6.1; rr(s,rx,ly,6.0,4.5,NAVY)
t(s,rx+0.3,ly+0.25,5.4,0.4,"LE PROJET",sz=14,bd=True,c=GOLD)
t(s,rx+0.3,ly+0.7,5.4,0.5,"Un ERP intelligent",sz=22,bd=True,c=WHITE)
t(s,rx+0.3,ly+1.25,5.4,3.0,
  "Concevoir et developper une application web de gestion integree qui digitalise "
  "les operations commerciales de Dind'Or et introduit une aide a la decision basee "
  "sur l'IA pour optimiser les achats et maximiser la marge.",sz=13,c=CREAM)
t(s,0.6,6.85,12,0.3,"DIGITALISATION  +  IA  =  AVANTAGE COMPETITIF",
  sz=11,bd=True,c=GOLD,a=PP_ALIGN.CENTER,it=True)
pg(s,3)

# ===== SLIDE 4 - OBJECTIFS =====
s=prs.slides.add_slide(B); bg(s,CREAM)
hdr(s,"02 | Objectifs","Ce que le projet doit accomplir")
rr(s,0.6,1.95,12.1,0.85,NAVY,0.1)
t(s,0.85,2.05,2,0.3,"OBJECTIF GENERAL",sz=10,bd=True,c=GOLD)
t(s,0.85,2.32,11.5,0.45,
  "Digitaliser les operations + integrer une IA d'aide a la decision sur les achats",
  sz=14,c=WHITE)
cards=[("Gestion commerciale","Achats, ventes, stock, etat financier en temps reel"),
       ("Module IA predictif","Prevision des prix a 7, 15 et 30 jours / Buy ou Hold"),
       ("Portail multi-roles","Acces dedies Admin, Client et Employe"),
       ("Securite & scalabilite","Auth JWT, architecture 3 couches modulaire")]
cw,ch=2.95,3.6
for i,(tt,tx) in enumerate(cards):
    cx=0.6+i*(cw+0.15); cy=3.0
    rr(s,cx,cy,cw,ch,WHITE,0.06); r(s,cx,cy,cw,0.08,GOLD)
    o(s,cx+0.3,cy+0.3,0.55,NAVY)
    t(s,cx+0.3,cy+0.42,0.55,0.45,f"0{i+1}",sz=16,bd=True,c=GOLD,a=PP_ALIGN.CENTER)
    t(s,cx+0.2,cy+1.05,cw-0.4,0.6,tt,sz=15,bd=True,c=DARK)
    t(s,cx+0.2,cy+1.7,cw-0.4,1.7,tx,sz=11,c=SLATE)
pg(s,4)

# ===== SLIDE 5 - PROBLEMATIQUE =====
s=prs.slides.add_slide(B); bg(s,CREAM)
hdr(s,"03 | Problematique","Le prix d'achat de la volaille est imprevisible")
lx,ly,lw=0.6,2.05,6.3
rr(s,lx,ly,lw,4.95,WHITE); r(s,lx,ly,lw,0.55,NAVY)
t(s,lx+0.3,ly+0.13,lw-0.5,0.35,"CAUSES DE L'INSTABILITE",sz=13,bd=True,c=GOLD)
ca=[("Ramadan & Aid","+15% a +30%"),("Ete (juillet-aout)","+10% a +15%"),
    ("Vacances scolaires","+5% a +8%"),("Marche hebdomadaire","Imprevisible")]
cy=ly+0.85
for f,p in ca:
    rr(s,lx+0.3,cy,lw-0.6,0.85,CREAM,0.1); r(s,lx+0.3,cy,0.08,0.85,GOLD)
    t(s,lx+0.55,cy+0.12,3.5,0.35,f,sz=13,bd=True,c=DARK)
    t(s,lx+0.55,cy+0.45,3.5,0.35,"Facteur saisonnier",sz=10,c=SLATE,it=True)
    rr(s,lx+lw-1.65,cy+0.18,1.25,0.5,RED,0.3)
    t(s,lx+lw-1.65,cy+0.25,1.25,0.4,p,sz=12,bd=True,c=WHITE,a=PP_ALIGN.CENTER)
    cy+=1.0
rx=lx+lw+0.3; rw=13.33-rx-0.6
rr(s,rx,ly,rw,2.3,NAVY)
t(s,rx+0.3,ly+0.2,rw-0.6,0.35,"CONSEQUENCE",sz=12,bd=True,c=GOLD)
t(s,rx+0.3,ly+0.6,rw-0.6,0.5,"Decisions a l'instinct",sz=20,bd=True,c=WHITE)
t(s,rx+0.3,ly+1.2,rw-0.6,1.0,
  "Achats au mauvais moment, ruptures ou surstock, aucun historique structure.",sz=12,c=CREAM)
rr(s,rx,ly+2.5,rw,2.45,WHITE); r(s,rx,ly+2.5,rw,0.08,RED)
t(s,rx+0.3,ly+2.7,rw-0.6,0.35,"IMPACT FINANCIER ESTIME",sz=11,bd=True,c=RED)
t(s,rx+0.3,ly+3.15,rw-0.6,0.9,"12 600 DT",sz=44,bd=True,c=DARK)
t(s,rx+0.3,ly+4.05,rw-0.6,0.4,"de marge perdue par an (~350 DT par commande)",
  sz=11,c=SLATE,it=True)
pg(s,5)

# ===== SLIDE 6 - SOLUTION =====
s=prs.slides.add_slide(B); bg(s,CREAM)
hdr(s,"04 | Solution proposee","Un ERP complet + une IA d'aide a la decision")
lx,ly,w=0.6,2.05,6.15
rr(s,lx,ly,w,4.95,WHITE); r(s,lx,ly,w,0.55,NAVY)
t(s,lx+0.3,ly+0.13,w-0.5,0.35,"PARTIE 1 - ERP CLASSIQUE",sz=13,bd=True,c=GOLD)
em=[("Achats","Factures fournisseurs, paiements multi-modes"),
    ("Ventes","Bons de livraison, facturation, encaissements"),
    ("Stock","Mouvements, seuil mini, PUMP"),
    ("RH","Fiches employes, pointage, conges, paie"),
    ("Portail","Admin / Client / Employe")]
cy=ly+0.85
for i,(tt,d) in enumerate(em):
    o(s,lx+0.3,cy+0.07,0.32,GOLD)
    t(s,lx+0.3,cy+0.13,0.32,0.28,str(i+1),sz=12,bd=True,c=NAVY,a=PP_ALIGN.CENTER)
    t(s,lx+0.75,cy,w-1,0.32,tt,sz=13,bd=True,c=DARK)
    t(s,lx+0.75,cy+0.32,w-1,0.32,d,sz=10,c=SLATE); cy+=0.78
rx=lx+w+0.3; rw=13.33-rx-0.6
rr(s,rx,ly,rw,4.95,NAVY)
t(s,rx+0.3,ly+0.25,rw-0.6,0.35,"PARTIE 2 - MODULE IA",sz=13,bd=True,c=GOLD)
t(s,rx+0.3,ly+0.7,rw-0.6,0.5,"Prediction des prix d'achat",sz=20,c=WHITE,bd=True)
ai=[("Modele","Gradient Boosting Regressor - 22 variables"),
    ("Horizons","Predictions a 7, 15 et 30 jours"),
    ("Saisonnalite","Calendrier tunisien : Ramadan, Aid, ete"),
    ("Recommandation","Buy maintenant / Hold (attendre)"),
    ("Validation","Walk-forward (MAE, MAPE, precision)")]
cy=ly+1.4
for k,v in ai:
    r(s,rx+0.3,cy+0.15,0.05,0.35,GOLD)
    t(s,rx+0.45,cy,rw-0.7,0.32,k,sz=11,bd=True,c=GOLD)
    t(s,rx+0.45,cy+0.32,rw-0.7,0.34,v,sz=11,c=CREAM); cy+=0.7
pg(s,6)

# ===== SLIDE 7 - ARCHITECTURE =====
s=prs.slides.add_slide(B); bg(s,CREAM)
hdr(s,"05 | Architecture technique","3 couches + microservice IA")
ly=[("PRESENTATION","Angular 19","SPA - Portails Admin/Client/Employe",GOLD),
    ("APPLICATION","Spring Boot 3","API REST - JWT - Logique metier",NAVY),
    ("DONNEES","MySQL","Persistance ERP - tables relationnelles",SLATE)]
bw,bh,gp=3.5,2.4,0.4
tw=bw*3+gp*2; sx=(W-tw)/2; yy=2.3
for i,(la,th,ds,cl) in enumerate(ly):
    x=sx+i*(bw+gp)
    rr(s,x,yy,bw,bh,WHITE); r(s,x,yy,bw,0.1,cl)
    t(s,x+0.2,yy+0.25,bw-0.4,0.35,la,sz=11,bd=True,c=cl)
    t(s,x+0.2,yy+0.65,bw-0.4,0.6,th,sz=22,bd=True,c=DARK)
    t(s,x+0.2,yy+1.35,bw-0.4,0.9,ds,sz=11,c=SLATE)
    if i<2:
        ax=x+bw+0.05; ay=yy+bh/2-0.15
        a=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(ax),Inches(ay),Inches(gp-0.1),Inches(0.3))
        a.fill.solid(); a.fill.fore_color.rgb=GOLD; a.line.fill.background(); a.shadow.inherit=False
ay=yy+bh+0.5; aw=8.0; ax=(W-aw)/2
rr(s,ax,ay,aw,1.6,NAVY,0.08); r(s,ax,ay,0.1,1.6,GOLD)
t(s,ax+0.4,ay+0.2,aw-0.6,0.35,"MICROSERVICE IA (decouple)",sz=11,bd=True,c=GOLD)
t(s,ax+0.4,ay+0.55,aw-0.6,0.5,"Python FastAPI",sz=22,bd=True,c=WHITE)
t(s,ax+0.4,ay+1.05,aw-0.6,0.5,
  "Endpoint /predict - previsions consommees par Spring Boot",sz=11,c=CREAM)
a2=s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,Inches(W/2-0.15),Inches(yy+bh+0.1),Inches(0.3),Inches(0.35))
a2.fill.solid(); a2.fill.fore_color.rgb=GOLD; a2.line.fill.background(); a2.shadow.inherit=False
pg(s,7)

# ===== SLIDE 8 - ENVIRONNEMENT =====
s=prs.slides.add_slide(B); bg(s,CREAM)
hdr(s,"06 | Environnement de travail","Outils et technologies")
cats=[("FRONTEND",[("Angular 19",GOLD),("TypeScript",SLATE),("Tailwind CSS",SLATE)]),
      ("BACKEND",[("Spring Boot 3",GOLD),("Java 17+",SLATE),("JWT Security",SLATE)]),
      ("DONNEES & IA",[("MySQL",GOLD),("Python FastAPI",GOLD),("scikit-learn",SLATE),("Pandas",SLATE)]),
      ("OUTILS",[("Git / GitHub",SLATE),("Postman",SLATE),("VS Code / IntelliJ",SLATE),("Docker",SLATE)])]
cw,ch=2.95,4.55; cy=2.1
for i,(ct,tx) in enumerate(cats):
    cx=0.6+i*(cw+0.15)
    rr(s,cx,cy,cw,ch,WHITE); r(s,cx,cy,cw,0.5,NAVY)
    t(s,cx+0.2,cy+0.12,cw-0.4,0.3,ct,sz=12,bd=True,c=GOLD,a=PP_ALIGN.CENTER)
    ty=cy+0.75
    for tc,cl in tx:
        rr(s,cx+0.3,ty,cw-0.6,0.55,CREAM,0.25); r(s,cx+0.3,ty,0.06,0.55,cl)
        t(s,cx+0.45,ty+0.13,cw-0.7,0.3,tc,sz=12,bd=True,c=DARK)
        ty+=0.7
pg(s,8)

# ===== SLIDE 9 - SCRUM =====
s=prs.slides.add_slide(B); bg(s,CREAM)
hdr(s,"07 | Methodologie SCRUM","Agile, iteratif, transparent")
rr(s,0.6,2.0,12.1,0.85,NAVY,0.1)
t(s,0.85,2.1,2,0.3,"POURQUOI SCRUM",sz=10,bd=True,c=GOLD)
t(s,0.85,2.35,11.5,0.5,
  "Livraisons incrementales, feedback rapide, adaptation continue ERP + IA",sz=13,c=WHITE)
lx,ly=0.6,3.05
rr(s,lx,ly,5.8,3.8,WHITE)
t(s,lx+0.3,ly+0.25,5.2,0.35,"ROLES & ARTEFACTS",sz=12,bd=True,c=NAVY)
ro=[("Product Owner","Encadrant + responsable Dind'Or"),
    ("Scrum Master","Etudiant - facilitateur"),
    ("Equipe Dev","Etudiant - dev fullstack + IA"),
    ("Product Backlog","User stories priorisees"),
    ("Sprint Backlog","Selection pour chaque sprint"),
    ("Rituels","Daily / Review / Retro")]
ry=ly+0.75
for k,d in ro:
    o(s,lx+0.3,ry+0.1,0.12,GOLD)
    t(s,lx+0.55,ry,2.4,0.35,k,sz=11,bd=True,c=DARK)
    t(s,lx+2.85,ry,2.85,0.35,d,sz=10,c=SLATE,it=True); ry+=0.47
rx=lx+6.1; rr(s,rx,ly,6.0,3.8,NAVY)
t(s,rx+0.3,ly+0.25,5.4,0.35,"DECOUPAGE EN SPRINTS",sz=12,bd=True,c=GOLD)
sp=[("Sprint 0","Cadrage, specs, maquettes"),
    ("Sprint 1","Auth JWT + module Achats"),
    ("Sprint 2","Ventes + Stock + Tresorerie"),
    ("Sprint 3","RH + Portail Client/Employe"),
    ("Sprint 4","Microservice IA + dashboard"),
    ("Sprint 5","Tests, deploiement, recette")]
sy=ly+0.75
for k,d in sp:
    r(s,rx+0.3,sy+0.13,0.05,0.28,GOLD)
    t(s,rx+0.45,sy,1.6,0.4,k,sz=12,bd=True,c=GOLD)
    t(s,rx+2.05,sy,3.8,0.4,d,sz=11,c=CREAM); sy+=0.47
pg(s,9)

# ===== SLIDE 10 - CONCLUSION =====
s=prs.slides.add_slide(B); bg(s,NAVY)
b2=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(-2),Inches(4.5),Inches(7),Inches(7))
b2.fill.solid(); b2.fill.fore_color.rgb=GOLD; b2.line.fill.background(); b2.shadow.inherit=False
sp2=b2.fill.fore_color._xFill
al2=etree.SubElement(sp2.find(qn('a:srgbClr')),qn('a:alpha')); al2.set('val','12000')
r(s,0.6,0.7,0.4,0.35,GOLD)
t(s,1.1,0.65,6,0.4,"08 | CONCLUSION",sz=12,bd=True,c=GOLD)
t(s,0.6,1.2,12,0.85,"Un projet a forte valeur ajoutee",sz=36,bd=True,c=WHITE)
mx,my=0.6,2.6; mw,mh=3.95,2.5
me=[("100%","Operations digitalisees","Achats, ventes, stock, RH"),
    ("7-30 j","Anticipation des prix","Predictions multi-horizons"),
    ("12K+ DT","Marge recuperee","Gain annuel estime")]
for i,(v,l,sb) in enumerate(me):
    x=mx+i*(mw+0.2)
    rr(s,x,my,mw,mh,NAVY_D); r(s,x,my,mw,0.08,GOLD)
    t(s,x+0.3,my+0.35,mw-0.6,0.95,v,sz=42,bd=True,c=GOLD)
    t(s,x+0.3,my+1.35,mw-0.6,0.5,l,sz=14,bd=True,c=WHITE)
    t(s,x+0.3,my+1.9,mw-0.6,0.5,sb,sz=10,c=CREAM,it=True)
r(s,0.6,5.45,0.5,0.05,GOLD)
t(s,0.6,5.55,12,0.35,"PERSPECTIVES",sz=11,bd=True,c=GOLD)
t(s,0.6,5.95,12.1,1.4,
  "Application mobile - CRM avance - Predictions multi-produits - BI temps reel - Paiement en ligne",
  sz=14,c=CREAM)
t(s,0.6,6.85,12.1,0.4,"MERCI POUR VOTRE ATTENTION - QUESTIONS ?",
  sz=14,bd=True,c=GOLD,a=PP_ALIGN.CENTER,it=True)

prs.save("Dindor_PFE_Presentation.pptx")
print("OK: Dindor_PFE_Presentation.pptx genere.")
